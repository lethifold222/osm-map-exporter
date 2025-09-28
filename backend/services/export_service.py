import asyncio
import os
import json
from pathlib import Path
from typing import Dict, Any, List
from datetime import datetime
import httpx
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

from models.schemas import ExportPDFRequest, ExportPPTXRequest

class ExportService:
    def __init__(self):
        self.export_dir = Path("exports")
        self.export_dir.mkdir(exist_ok=True)
        self.template_dir = Path("export/templates")
        self.template_dir.mkdir(parents=True, exist_ok=True)
    
    async def export_pdf(self, request: ExportPDFRequest) -> str:
        """Export map data as PDF using Playwright"""
        filename = f"context_maps_{self._sanitize_filename(request.title)}_{datetime.now().strftime('%Y%m%d')}.pdf"
        file_path = self.export_dir / filename
        
        # Create HTML template for PDF export
        html_content = await self._create_pdf_html(request)
        
        # Save HTML template
        template_path = self.template_dir / f"pdf_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        with open(template_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        # Generate PDF using Playwright
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()
            
            # Load the HTML template
            await page.goto(f"file://{template_path.absolute()}")
            
            # Wait for map to load
            await page.wait_for_timeout(3000)
            
            # Generate PDF
            await page.pdf(
                path=str(file_path),
                format='A4',
                print_background=True,
                margin={'top': '0.5in', 'bottom': '0.5in', 'left': '0.5in', 'right': '0.5in'}
            )
            
            await browser.close()
        
        # Clean up template
        template_path.unlink()
        
        return str(file_path)
    
    async def export_pptx(self, request: ExportPPTXRequest) -> str:
        """Export map data as PowerPoint presentation"""
        filename = f"context_maps_{self._sanitize_filename(request.title)}_{datetime.now().strftime('%Y%m%d')}.pptx"
        file_path = self.export_dir / filename
        
        # Create new presentation
        prs = Presentation()
        
        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = request.title
        if request.subtitle:
            title_slide.placeholders[1].text = request.subtitle
        
        # Add metadata slide
        metadata_slide = prs.slides.add_slide(prs.slide_layouts[1])
        metadata_slide.shapes.title.text = "Տվյալների մասին / Data Information"
        
        # Add content for each layer
        for layer in request.layers:
            if layer == "roads":
                await self._add_roads_slide(prs, request)
            elif layer == "buildings":
                await self._add_buildings_slide(prs, request)
            elif layer == "amenities":
                await self._add_amenities_slide(prs, request)
            elif layer == "pois":
                await self._add_pois_slide(prs, request)
        
        # Save presentation
        prs.save(str(file_path))
        
        return str(file_path)
    
    async def _create_pdf_html(self, request: ExportPDFRequest) -> str:
        """Create HTML template for PDF export"""
        # This would create a complete HTML page with embedded map
        # For now, return a basic template
        return f"""
        <!DOCTYPE html>
        <html lang="hy">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{request.title}</title>
            <style>
                body {{
                    font-family: 'Arial', sans-serif;
                    margin: 0;
                    padding: 20px;
                    background: white;
                }}
                .header {{
                    text-align: center;
                    margin-bottom: 30px;
                    border-bottom: 2px solid #333;
                    padding-bottom: 20px;
                }}
                .map-container {{
                    width: 100%;
                    height: 600px;
                    border: 1px solid #ccc;
                    margin-bottom: 20px;
                }}
                .legend {{
                    display: flex;
                    flex-wrap: wrap;
                    gap: 20px;
                    margin-top: 20px;
                }}
                .legend-item {{
                    display: flex;
                    align-items: center;
                    gap: 10px;
                }}
                .legend-color {{
                    width: 20px;
                    height: 20px;
                    border: 1px solid #333;
                }}
                .footer {{
                    margin-top: 30px;
                    font-size: 12px;
                    color: #666;
                    text-align: center;
                }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>{request.title}</h1>
                {f'<h2>{request.subtitle}</h2>' if request.subtitle else ''}
                <p>Ամսաթիվ / Date: {datetime.now().strftime('%Y-%m-%d')}</p>
            </div>
            
            <div class="map-container" id="map">
                <!-- Map will be rendered here -->
                <p style="text-align: center; margin-top: 250px;">Քարտեզը բեռնվում է... / Map loading...</p>
            </div>
            
            <div class="legend">
                <div class="legend-item">
                    <div class="legend-color" style="background-color: #ff0000;"></div>
                    <span>Ճանապարհներ / Roads</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background-color: #0000ff;"></div>
                    <span>Շենքեր / Buildings</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background-color: #00ff00;"></div>
                    <span>Հարմարություններ / Amenities</span>
                </div>
            </div>
            
            <div class="footer">
                <p>Տվյալներ OpenStreetMap-ից / Data from OpenStreetMap</p>
                <p>Գեներացվել է OSM Map Exporter-ով / Generated by OSM Map Exporter</p>
            </div>
        </body>
        </html>
        """
    
    async def _add_roads_slide(self, prs: Presentation, request: ExportPPTXRequest):
        """Add roads layer slide to presentation"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Ճանապարհային ցանց / Road Network"
        
        # Add map image placeholder
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        
        slide.shapes.add_picture("placeholder_roads.png", left, top, width, height)
    
    async def _add_buildings_slide(self, prs: Presentation, request: ExportPPTXRequest):
        """Add buildings layer slide to presentation"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Շենքեր / Buildings"
        
        # Add map image placeholder
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        
        slide.shapes.add_picture("placeholder_buildings.png", left, top, width, height)
    
    async def _add_amenities_slide(self, prs: Presentation, request: ExportPPTXRequest):
        """Add amenities layer slide to presentation"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Հարմարություններ / Amenities"
        
        # Add map image placeholder
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        
        slide.shapes.add_picture("placeholder_amenities.png", left, top, width, height)
    
    async def _add_pois_slide(self, prs: Presentation, request: ExportPPTXRequest):
        """Add POIs layer slide to presentation"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Հետաքրքրության կետեր / Points of Interest"
        
        # Add map image placeholder
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        
        slide.shapes.add_picture("placeholder_pois.png", left, top, width, height)
    
    def _sanitize_filename(self, filename: str) -> str:
        """Sanitize filename for safe file system usage"""
        import re
        # Remove or replace invalid characters
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # Limit length
        return filename[:50]
